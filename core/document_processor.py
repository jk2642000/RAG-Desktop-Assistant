import os
import PyPDF2
from docx import Document
from typing import List, Dict
import hashlib
import csv
from .logger import rag_logger

try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class DocumentProcessor:
    def __init__(self):
        self.supported_formats = ['.pdf', '.txt', '.docx', '.csv']
        if EXCEL_AVAILABLE:
            self.supported_formats.extend(['.xlsx', '.xls'])
        if PPTX_AVAILABLE:
            self.supported_formats.append('.pptx')
    
    def process_document(self, file_path: str) -> Dict:
        """Process a single document and return structured data"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        # Extract text based on file type
        if file_ext == '.pdf':
            text = self._extract_pdf_text(file_path)
        elif file_ext == '.txt':
            text = self._extract_txt_text(file_path)
        elif file_ext == '.docx':
            text = self._extract_docx_text(file_path)
        elif file_ext == '.csv':
            text = self._extract_csv_text(file_path)
        elif file_ext in ['.xlsx', '.xls'] and EXCEL_AVAILABLE:
            text = self._extract_excel_text(file_path)
        elif file_ext == '.pptx' and PPTX_AVAILABLE:
            text = self._extract_pptx_text(file_path)
        else:
            raise ValueError(f"Parser not available for format: {file_ext}")
        
        # Create document chunks
        chunks = self._create_chunks(text)
        
        # Calculate metrics
        file_size = os.path.getsize(file_path)
        avg_chunk_size = sum(len(chunk) for chunk in chunks) / len(chunks) if chunks else 0
        file_hash = self._get_file_hash(file_path)
        
        # Log document processing metrics
        rag_logger.log_document_processed(
            doc_id=file_hash,
            filename=os.path.basename(file_path),
            file_size=file_size,
            chunk_count=len(chunks),
            avg_chunk_size=avg_chunk_size
        )
        
        return {
            'filename': os.path.basename(file_path),
            'filepath': file_path,
            'text': text,
            'chunks': chunks,
            'file_hash': file_hash,
            'chunk_count': len(chunks),
            'file_size': file_size,
            'avg_chunk_size': avg_chunk_size
        }
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"Error reading PDF: {str(e)}")
        return text.strip()
    
    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read()
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error reading DOCX: {str(e)}")
    
    def _create_chunks(self, text: str, chunk_size: int = 800, overlap: int = 150) -> List[str]:
        """Split text into overlapping chunks"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings near the chunk boundary
                for i in range(end, max(start + chunk_size - 100, start), -1):
                    if text[i] in '.!?':
                        end = i + 1
                        break
            
            chunk = text[start:end].strip()
            # Clean up extra whitespace and newlines
            chunk = ' '.join(chunk.split())
            if chunk and len(chunk) > 50:  # Skip very short chunks
                chunks.append(chunk)
            
            start = end - overlap
            if start >= len(text):
                break
        
        return chunks
    
    def _extract_csv_text(self, file_path: str) -> str:
        """Extract text from CSV file"""
        try:
            text_parts = []
            with open(file_path, 'r', encoding='utf-8', newline='') as file:
                # Try to detect delimiter
                sample = file.read(1024)
                file.seek(0)
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter
                
                reader = csv.reader(file, delimiter=delimiter)
                headers = next(reader, None)
                
                if headers:
                    text_parts.append(f"CSV Headers: {', '.join(headers)}")
                
                for row_num, row in enumerate(reader, 1):
                    if row_num > 1000:  # Limit rows to prevent huge files
                        text_parts.append(f"... (truncated after 1000 rows)")
                        break
                    row_text = ' | '.join(str(cell) for cell in row if cell)
                    if row_text.strip():
                        text_parts.append(f"Row {row_num}: {row_text}")
                
            return '\n'.join(text_parts)
        except Exception as e:
            raise Exception(f"Error reading CSV: {str(e)}")
    
    def _extract_excel_text(self, file_path: str) -> str:
        """Extract text from Excel file"""
        try:
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"\n=== Sheet: {sheet_name} ===")
                
                # Get headers from first row
                headers = []
                for cell in sheet[1]:
                    if cell.value:
                        headers.append(str(cell.value))
                
                if headers:
                    text_parts.append(f"Headers: {', '.join(headers)}")
                
                # Process data rows (limit to 1000 rows)
                for row_num, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), 2):
                    if row_num > 1000:
                        text_parts.append("... (truncated after 1000 rows)")
                        break
                    
                    row_data = [str(cell) for cell in row if cell is not None]
                    if row_data:
                        text_parts.append(f"Row {row_num}: {' | '.join(row_data)}")
            
            workbook.close()
            return '\n'.join(text_parts)
        except Exception as e:
            raise Exception(f"Error reading Excel: {str(e)}")
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            presentation = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_parts.append(f"\n=== Slide {slide_num} ===")
                
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if hasattr(shape, 'table'):
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(' | '.join(row_text))
                
                if slide_text:
                    text_parts.extend(slide_text)
                else:
                    text_parts.append("(No text content)")
            
            return '\n'.join(text_parts)
        except Exception as e:
            raise Exception(f"Error reading PowerPoint: {str(e)}")
    
    def _get_file_hash(self, file_path: str) -> str:
        """Generate hash for file to detect changes"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()